import os
import csv
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_AUTO_SIZE
import argparse
import subprocess
import base64
import requests
from concurrent.futures import ThreadPoolExecutor
import subprocess

executor = ThreadPoolExecutor(max_workers=4)  # Customize thread count per system

def ensure_ollama_llava_running():
    try:
        subprocess.Popen(["ollama", "run", "llava"], creationflags=subprocess.CREATE_NEW_CONSOLE)
    except Exception as e:
        print(f"[ERROR] Couldn't launch Ollama LLaVA: {e}")

def convert_ppt_to_pptx(filepath):
    soffice = r"C:\\Program Files\\LibreOffice\\program\\soffice.exe"
    input_dir = os.path.dirname(filepath)

    subprocess.run([
        soffice,
        "--headless",
        "--convert-to", "pptx",
        filepath,
        "--outdir", input_dir
    ], check=True)

class PowerPointExtractor:
    def __init__(self, ppt_path, session_dir):
        self.ppt_path = ppt_path
        self.session_dir = session_dir
        self.image_output_dir = os.path.join(session_dir, "images")
        os.makedirs(self.image_output_dir, exist_ok=True)
        self.cur_image_index = 0
        self.invalid_images = []

    def save_image(self, image, name):
        image_bytes = image.blob
        name = name + f'_{self.cur_image_index}.{image.ext}'
        full_path = os.path.join(self.image_output_dir, os.path.basename(name))
        print(full_path) #printing to showcase multithreading
        with open(full_path, 'wb') as f:
            f.write(image_bytes)
        self.cur_image_index += 1
        return os.path.basename(full_path)

    def drill_for_images(self, shape, slide_idx, name):
        image_tuples = []
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for s in shape.shapes:
                image_tuples.extend(self.drill_for_images(s, slide_idx, name))
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                saved_image = self.save_image(shape.image, name)
                image_tuples.append((saved_image, shape))
            except:
                print(f'Could not process image {shape.name} on slide {slide_idx}.')
                self.invalid_images.append(f'Slide {slide_idx}: {shape.name}')
                image_tuples.append((f'INVALID: {shape.name}', None))
        else:
            try:
                if hasattr(shape, 'image'):
                    saved_image = self.save_image(shape.image, name)
                    image_tuples.append((saved_image, shape))
            except:
                pass
        return image_tuples 

    def get_slide_text(self, slide):
        text = "\n".join([
            shape.text.strip() for shape in slide.shapes
            if hasattr(shape, "text") and shape.text.strip()
        ])
        return "\n".join([line for line in text.splitlines() if line.strip()])

    def get_context_text(self, slides, current_index, window=1):
        main_slide_text = self.get_slide_text(slides[current_index])
        context_parts = [f"Slide Content:\n{main_slide_text.strip()}"]

        total = len(slides)
        neighbor_texts = []
        for offset in range(-window, window + 1):
            idx = current_index + offset
            if idx == current_index or not (0 <= idx < total):
                continue
            neighbor_slide_text = self.get_slide_text(slides[idx])
            if neighbor_slide_text.strip():
                neighbor_texts.append(neighbor_slide_text)

        if neighbor_texts:
            context_parts.append("\n\nRelated Slide Hints:\n" + "\n\n---\n\n".join(neighbor_texts))

        return "\n\n".join(context_parts)

    def generate_llava_caption(self, image_path, context_text):
        with open(image_path, "rb") as img_file:
            image_b64 = base64.b64encode(img_file.read()).decode()

        prompt = (
            "You are generating a short caption for an image on a presentation slide.\n"
            "Focus mainly on the **Slide Content** section.\n"
            "Only refer to **Related Slide Hints** if the Slide Content is vague or missing.\n"
            "The caption must be concise and strictly no more than 25 words.\n\n"
            f"{context_text.strip()}\n\nCaption:"
        )

        response = requests.post(
            "http://localhost:11434/api/generate",
            json={
                "model": "llava",
                "prompt": prompt,
                "images": [image_b64],
                "stream": False
            }
        )
        result = response.json()
        return result.get("response", "[No caption generated]")

    def add_caption_to_slide(self, slide, caption_text, image_shape):
        left = image_shape.left
        top = image_shape.top + image_shape.height + Inches(0.1)
        width = image_shape.width
        height = Inches(0.4)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.NONE

        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = caption_text

        font = run.font
        font.size = Pt(10)
        font.italic = True
        font.color.rgb = RGBColor(0, 0, 0)

        fill = textbox.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)

        line = textbox.line
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(200, 200, 200)

    def process_file(self):
        name_base = os.path.splitext(os.path.basename(self.ppt_path))[0]
        ppt = Presentation(self.ppt_path)
        out_csv_path = os.path.join(self.session_dir, f'{name_base}_captions.csv')

        TOP_MARGIN_THRESHOLD = 1 * 360000

        with open(out_csv_path, 'w', encoding="utf-8", newline='') as file:
            writer = csv.writer(file)
            writer.writerow(["Slide", "Context", "Image", "Caption"])

            for i, slide in enumerate(ppt.slides):
                if i == 0:
                    continue

                slide_context = self.get_context_text(ppt.slides, i, window=1)
                image_name_part = os.path.join(self.image_output_dir, f'{name_base}_slide{i+1}')

                image_tuples = []
                for shape in slide.shapes:
                    # print(f"Image position: {shape.top / 360000:.2f} cm")
                    image_tuples.extend(self.drill_for_images(shape, i + 1, image_name_part))

                for image_file, shape in image_tuples:
                    if image_file.startswith('INVALID') or shape is None:
                        continue

                    if shape.top < TOP_MARGIN_THRESHOLD:
                        print(f"⚠ Skipping image {image_file} on slide {i+1} (too close to top)")
                        continue

                    full_image_path = os.path.join(self.image_output_dir, image_file)
                    caption = self.generate_llava_caption(full_image_path, slide_context)
                    self.add_caption_to_slide(slide, caption, shape)
                    writer.writerow([i + 1, slide_context, image_file, caption])

        output_ppt = os.path.join(self.session_dir, f'{name_base}_captioned.pptx')
        ppt.save(output_ppt)
        print(f"✅ Saved modified presentation as {output_ppt}")

        if self.invalid_images:
            print(f'⚠ WARNING: {len(self.invalid_images)} invalid images found: {self.invalid_images}')

        return output_ppt

def run_captioning_threaded(input_path, session_dir):
    future = executor.submit(PowerPointExtractor(input_path, session_dir).process_file)
    return future

def main():
    parser = argparse.ArgumentParser(description="LLaVA-powered image captioning for PowerPoint")
    parser.add_argument('--ppt', required=True, help='Path to the input .ppt or .pptx file')
    parser.add_argument('--out', default='session_data', help='Directory to save session data')
    args = parser.parse_args()

    input_path = args.ppt
    session_dir = args.out

    ensure_ollama_llava_running() # ensures that LLaVa is runnign already if not it will run it in new terminal

    if not os.path.exists(input_path):
        print(f"[❌] File not found: {input_path}")
        return

    if input_path.lower().endswith('.ppt'):
        convert_ppt_to_pptx(input_path)
        input_path = input_path.rsplit('.', 1)[0] + '.pptx'

    output_ppt = run_captioning_threaded(input_path, session_dir).result()
    print(f"[✅] Processing complete: {output_ppt}")

if __name__ == '__main__':
    main()
